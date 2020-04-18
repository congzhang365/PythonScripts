from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
import re
import pandas as pd


def find_bold(str, tf):
    bold = re.compile('(<.*?>)')
    str_word = str.split(" ")
    # print(str_word)
    p = tf.add_paragraph()
    for i in range(len(str_word)):
        if not bold.search(str_word[i]):
            run = p.add_run()
            run.text = " " + str_word[i]
            font = run.font
            font.bold = False
            font.size = Pt(26)
        else:
            run = p.add_run()
            run.text = " " + str_word[i].replace('<', "").replace('>', "")
            font = run.font
            font.bold = True
            font.size = Pt(26)


def find_bracket(str,tf):
    brackets = re.compile('(\(.*?\))')
    # print(str_word)
    p = tf.add_paragraph()
    ## if yes, contents in brackets are grey and in italics
    if brackets.search(str):
        bracket_text = brackets.findall(str)[0]
        text_b_main = str.replace(bracket_text, '')
        text_b_bracket = bracket_text

        # configure 2nd text box


        run = p.add_run()
        run.text = '- ' + text_b_main
        font = run.font
        font.italic = False
        font.size = Pt(26)

        run2 = p.add_run()
        run2.text = text_b_bracket
        font = run2.font
        font.italic = True
        font.size = Pt(26)
        font.color.rgb = RGBColor(128, 128, 128)
    else:

        run2 = p.add_run()
        run2.text = '- ' + str
        font = run2.font
        font.size = Pt(26)


def find_bold_bracket(str, tf):
    bold = re.compile('(<.*?>)')
    brackets = re.compile('(\(.*?\))')
    both = re.compile('(<.*?>)|(\(.*?\))')
    run_text_list = list(filter(None, re.split(both, str)))
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = '- '
    for i in range(len(run_text_list)):

        if bold.search(run_text_list[i]):
            run = p.add_run()
            run.text = run_text_list[i].replace('<', "").replace('>', "")
            font = run.font
            font.bold = True
            font.size = Pt(26)

        elif brackets.search(run_text_list[i]):
            run = p.add_run()
            run.text = run_text_list[i]
            font = run.font
            font.italic = True
            font.size = Pt(26)
            font.color.rgb = RGBColor(128, 128, 128)
        else:
            run = p.add_run()
            run.text = " " + run_text_list[i]
            font = run.font
            font.bold = False
            font.size = Pt(26)


def create_slide(md_text, md_index):
    '''
    md_text is the entire mini dialogue
    md_index is the index number of the current mini dialogue
    '''

    # create a presentation using the blank template
    # prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    md_by_line = md_text.split('\n')
    text_a = md_by_line[0]
    text_b = md_by_line[1]

    if len(md_by_line) == 2:
        # configure 1st text box position and size (corresponding to the params in pptx)
        txBox1 = slide.shapes.add_textbox(left=Cm(1.91),
                                          top=Cm(5.92),
                                          width=Cm(21.59),
                                          height=Cm(2))
        tf1 = txBox1.text_frame
        tf1.word_wrap = True
        p = tf1.add_paragraph()
        p.text = '- ' + text_a
        p.font.italic = True
        p.font.size = Pt(26)
        p.font.color.rgb = RGBColor(0, 0, 128)


        # configure 2nd text box
        txBox2 = slide.shapes.add_textbox(left=Cm(1.91),
                                          top=Cm(8.92),
                                          width=Cm(21.59),
                                          height=Cm(6))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        find_bold_bracket(text_b, tf2)

    elif len(md_by_line) == 3:
        text_c = md_by_line[2]
        # configure 1st text box position and size (corresponding to the params in pptx)
        txBox1 = slide.shapes.add_textbox(left=Cm(1.91),
                                          top=Cm(5.92),
                                          width=Cm(21.59),
                                          height=Cm(2))
        tf1 = txBox1.text_frame
        tf1.word_wrap = True
        p = tf1.add_paragraph()
        p.text = '- ' + text_a
        p.font.italic = True
        p.font.size = Pt(26)
        p.font.color.rgb = RGBColor(0, 0, 128)

        # configure 2nd text box
        txBox2 = slide.shapes.add_textbox(left=Cm(1.91),
                                          top=Cm(8.92),
                                          width=Cm(21.59),
                                          height=Cm(6))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        find_bold_bracket(text_b, tf2)

        # configure 3rd text box
        txBox3 = slide.shapes.add_textbox(left=Cm(1.91),
                                          top=Cm(11.92),
                                          width=Cm(21.59),
                                          height=Cm(2))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p = tf3.add_paragraph()
        p.text = '- ' + text_c
        p.font.italic = True
        p.font.size = Pt(26)
        p.font.color.rgb = RGBColor(0, 0, 128)


    # create a text box at the bottom right for the index numbers
    txBox_index = slide.shapes.add_textbox(left=Cm(22.77),
                                      top=Cm(17.66),
                                      width=Cm(1.22),
                                      height=Cm(1.01))
    tf3 = txBox_index.text_frame
    p = tf3.add_paragraph()
    p.text = md_index  # mini_dialogue index
    p.font.size = Pt(12)


def create_slide_star(text_a, text_b, md_index):
    '''
    This function is to mask the first turn as stars
    '''

    # create a presentation using the blank template
    # prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # configure 1st text box position and size (corresponding to the params in pptx)
    txBox1 = slide.shapes.add_textbox(left=Cm(1.91),
                                      top=Cm(5.92),
                                      width=Cm(21.59),
                                      height=Cm(2))
    tf1 = txBox1.text_frame
    tf1.word_wrap = True
    p = tf1.add_paragraph()
    square_brackets = re.compile('\[(.*?)\]')

    ## if yes, contents in brackets are replaced with '***'
    if square_brackets.search(text_a):
        square_brackets_text = square_brackets.findall(text_a)[0]
        text_a_replace = text_a.replace(square_brackets_text, '***')
    p.text = '- ' + text_a_replace

    # set font styles (italic, size 28, grey @SPRINT)
    p.font.italic = True
    p.font.size = Pt(26)
    p.font.color.rgb = RGBColor(0, 0, 128)

    # check if the text_b has any contents in brackets:
    brackets = re.compile('(\(.*?\))')

    ## if yes, contents in brackets are grey and in italics
    if brackets.search(text_b):
        bracket_text = brackets.findall(text_b)[0]
        text_b_main = text_b.replace(bracket_text, '')
        text_b_bracket = bracket_text

        # configure 2nd text box
        txBox2 = slide.shapes.add_textbox(left=Cm(1.91),
                                          top=Cm(7.92),
                                          width=Cm(21.59),
                                          height=Cm(6))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p = tf2.add_paragraph()
        run = p.add_run()
        run.text = '- ' + text_b_main
        font = run.font
        font.italic = False
        font.size = Pt(26)

        run2 = p.add_run()
        run2.text = text_b_bracket
        font = run2.font
        font.italic = True
        font.size = Pt(26)
        font.color.rgb = RGBColor(128, 128, 128)
    else:
        txBox2 = slide.shapes.add_textbox(left=Cm(1.91),
                                          top=Cm(8.92),
                                          width=Cm(21.59),
                                          height=Cm(6))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p = tf2.add_paragraph()
        run2 = p.add_run()
        run2.text = '- ' + text_b
        font = run2.font
        font.size = Pt(26)


    # create a text box at the bottom right for the index numbers
    txBox3 = slide.shapes.add_textbox(left=Cm(22.77),
                                      top=Cm(17.66),
                                      width=Cm(1.22),
                                      height=Cm(1.01))
    tf3 = txBox3.text_frame
    p = tf3.add_paragraph()
    p.text = md_index  # mini_dialogue index
    p.font.size = Pt(12)


if __name__ == '__main__':
    input = "C:/Users/sprin/SPRINT Dropbox/Academic Research/Production_Study/- mini-dialogues/mini-dialogues_GR/Rand_Formula.xlsx"
    # input = "C:/Users/sprin/Desktop/test/mini.txt"
    # df = pd.read_table(input). sort_values(by=['RAND1'])    # use this line to read txt
    df = pd.read_excel(input, "all").sort_values(by=['RAND1']) # sort df by Column 'RAND1'
    index_list = df.iloc[:, 0].to_list()
    text_list = df.iloc[:, 1].to_list()

    prs = Presentation()

    for i in range(0, len(df)):
        md_index = index_list[i]
        md_text = text_list[i]
        create_slide(md_text, md_index)

    # save your pptx file.
    # remember to change randomised version number.
    prs.save("C:/Users/sprin/SPRINT Dropbox/Academic Research/Production_Study/- mini-dialogues/mini-dialogues_GR/test.pptx")


'''
This script is for automatically formatting two-turn mini dialogues from txt to pptx
    To run this file, you need to install
        python-pptx
        regex
        pandas
    input file: tab-delimited txt file or excel file. Comment out the option that you don't need.
    create_slide: create a slide with three text boxes:     
    create_slide_star: Lines in text_a are presented as - [***].

Dr Cong Zhang 30/03/2020 @SPRINT
Last commit: 01/04/2020 
'''
